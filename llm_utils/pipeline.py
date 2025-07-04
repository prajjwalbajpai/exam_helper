from pptx import Presentation
from docx import Document
import fitz  # PyMuPDF
from langchain_text_splitters import RecursiveCharacterTextSplitter
import re
from langgraph.store.memory import InMemoryStore
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_groq import ChatGroq
import os
from dotenv import load_dotenv

load_dotenv()
GROQ_API_KEY = os.getenv("GROQ_API_KEY")


def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    text = []
    for para in doc.paragraphs:
        text.append(para.text)
    return "\n".join(text)

def extract_text_from_pdf(file_path):
    doc = fitz.open(file_path)
    text = []
    for page in doc:
        text.append(page.get_text())
    return "\n".join(text)

def split_text_into_chunks(text, chunk_size=1000, chunk_overlap=100):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = text_splitter.split_text(text)
    return chunks

def clean_text(text):
    text = re.sub(r'\n+', '\n', text)
    text = re.sub(r'[^\x00-\x7F]+', ' ', text)
    text = re.sub(r'Page \d+ of \d+', '', text)
    text = re.sub(r'Thank You', '', text)
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

def get_uploaded_file_paths():
    upload_dir = os.path.join(os.path.dirname(__file__), '..', 'uploads')
    upload_dir = os.path.abspath(upload_dir)  

    file_paths = [
        os.path.join(upload_dir, fname)
        for fname in os.listdir(upload_dir)
        if os.path.isfile(os.path.join(upload_dir, fname))
    ]

    return file_paths

def files_preprocess():
    files = get_uploaded_file_paths()
    file_chunks = []
    for file in files:
        # print(f"Processing file: {file}")
        ext = os.path.splitext(file)[1].lower().lstrip('.')
        file_text = ""

        if ext == "docx":
            file_text += extract_text_from_docx(file)
        elif ext == "pdf":
            file_text += extract_text_from_pdf(file)
        elif ext == "ppt" or ext == "pptx":
            file_text += extract_text_from_pptx(file)
            
        improved_file = clean_text(file_text)
        file_chunks += split_text_into_chunks(improved_file)
    
    return file_chunks        

def embed_data(file_chunks):
    embedding_model = HuggingFaceEmbeddings(
        model_name="BAAI/bge-base-en",
        model_kwargs={"device": "cuda"} 
    )

    store = InMemoryStore(
        index={
            "dims": 1536,
            "embed": embedding_model.embed_documents  
        }
    )

    for i, chunk in enumerate(file_chunks):
        store.put(("docs",), f"chunk_{i}", {"text": chunk})
    
    return store

def get_result_llm(prompt):
    llm = ChatGroq(
        model="llama-3.1-8b-instant",
        temperature=0.0,
        max_retries=2,
        api_key=GROQ_API_KEY
    )
    return llm.invoke(prompt).content


def get_answers(question):
    chunks = files_preprocess()
    store = embed_data(chunks)

    results = store.search(("docs",), query=question , limit=5)
    context = "\n\n".join([doc.value["text"] for doc in results])
    prompt = f"""You are an AI assistant. Use the context below to answer the question.

    Context:
    {context}

    Question: {question}

    Answer:"""

    return get_result_llm(prompt)

def get_imp_topics(answers):
    final_prompt = f"""
You will be given a long string of answers to previous years of question paper. Identify the recurring themes, key concepts, or main topics in these responses. 
Output a prioritized list of the most exam-relevant topics, sorted by importance or frequency (most frequent first). 
Provide only the list of the top 3 to 10 topics, with no extra summary.

Answers: {answers}

"""
    return get_result_llm(final_prompt)



